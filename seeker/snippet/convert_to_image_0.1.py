#date: 2022-10-27T17:13:47Z
#url: https://api.github.com/gists/674fbee991e4bea7cc8bbfee9b4701ae
#owner: https://api.github.com/users/kokotu0

import tkinter.filedialog
import pandas as pd
import pptx
import shutil
from pack_set import *
import openpyxl

def organize_foler_select():
    organize_folder_ent['state']='normal'
    organize_folder_ent.delete(0,tk.END)
    organize_directory=tkinter.filedialog.askdirectory(initialdir='/',title='폴더 선택')
    organize_folder_ent.insert(tk.END,organize_directory)
    organize_folder_ent['state']='readonly'
def open_files():
    overlap=0
    result=tkinter.filedialog.askopenfilenames(initialdir='/',title="파일 선택",filetypes=(("pptx files","*.pptx"),("all files","*.*")))
    #중복 자동제거
    result=list(set(result))
    for file_name in result:
        if file_name in file_list.get(0,tk.END):
            overlap+=1
            continue
        else:
            file_list.insert(tk.END,file_name)
    if overlap!=0 : tkinter.messagebox.showinfo(title="중복 확인",message="중복되는 항목 {}개를 제외하고 추가하였습니다.".format(overlap))
def delete_list():
    file_list.delete(file_list.curselection())
def delete_all():
    count=file_list.size()
    for i in range(count):
        file_list.delete(tk.END)
def save_file():
    pass
    # save_file_name=tkinter.filedialog.asksaveasfilename(initialdir='/',defaultextension=".xlsx",filetypes=(("xslx files","*.xlsx"),("all files","*.*")))
    # print(save_file_name)
    directory_set=file_list.get(0,tk.END)
    print(directory_set)
    # merge_and_release(directory_set,save_file_name)
    # tkinter.messagebox.showinfo(message="작업 완료")
def createfolder(directory):
    if not os.path.exists(directory):
        os.makedirs(directory)
def file_move():
    pass
    #pptx파일옮기기

def conver_to_image_ppts(file_names,organize_folder):
    APPLICATION = win32com.client.Dispatch("PowerPoint.Application")
    def slide_information():
        #날짜는 0슬라이드의 날짜에서 가져옴
        #이후 슬라이드 병합 프로그램에서 사용할 xl파일만들기
        wb=openpyxl.workbook.Workbook()
        df=pd.DataFrame()
        df.columns=['슬라이드 번호','설명','자료 출처 등 비고']

    def convert_to_image_ppt(file_name):
        prs=pptx.Presentation(file_name)
        PRESENTATION = APPLICATION.Presentations.Open(file_name, ReadOnly=False)
        foldername=organize_folder+file_name[file_name.rfind('/'):file_name.rfind('.')]
        createfolder(foldername)
        #pptx파일 옮기기
        shutil.copy(file_name, foldername)

        df = pd.DataFrame(columns=['슬라이드 번호', '설명', '자료 출처 등 비고'])
        for index, slide in enumerate(prs.slides):

            # print(foldername + '/' + str(index) + ".jpg")
            PRESENTATION.Slides[index].Export(str(foldername + '/' + str(index) + ".jpg").replace('/', '\\'), "JPG")
            try:
                title = slide.shapes.title.text
            except:
                title = ""

            x = pd.Series([index, title, ''], index=df.columns)
            df = pd.concat([df, x.to_frame().T], axis=0, ignore_index=True)
            #이제 df를 바탕으로 openpyxl을 이용하여 worksheet만들어서 저장하기
            
        df.to_excel(foldername+'/'+file_name[file_name.rfind('/'):file_name.rfind('.')]+'_슬라이드정리.xlsx', index=False)

    for file_name in file_names:
        convert_to_image_ppt(file_name)
    APPLICATION.Quit()
def make_excel():pass
# <editor-fold desc="tkinter 레이아웃">

root=tk.Tk()
root.title("이미지 변환 어플")

file_frame_lbl=tk.Label(root,text="PPT의 이미지 분리 및 폴더로 정리합니다. /n XLSX파일로 설명을 추가하거나 수정할 수 있습니다.")
file_frame_lbl.grid(row=0,column=0,columnspan=2)
organize_folder_ent=tk.Entry(root,state='readonly')

organize_folder_ent.grid(row=1,column=0,sticky='we')

organize_folder_btn=tk.Button(root,text="정리 폴더 선택",command=organize_foler_select)
organize_folder_btn.grid(row=1,column=1)

file_frame=tk.Frame(root,relief='groove',width=350,height=600)
file_frame.grid(row=2,column=0)
file_list=tk.Listbox(file_frame,width=80,height=50)
file_list.grid()
button_font=tkinter.font.Font(size=15)
button_frame=tk.Frame(root,relief='ridge',width=100,height=600)
button_frame.grid(row=2,column=1)


button_load=tk.Button(button_frame,text="불러오기",command=open_files)
button_load.grid(row=0,column=0,sticky='we')
button_delete=tk.Button(button_frame,text="선택 제거",command=delete_list)
button_delete.grid(row=1,column=0,sticky='we')
button_delete_all=tk.Button(button_frame,text="전체 제거",command=delete_all)
button_delete_all.grid(row=3,column=0,sticky='we')

button_save=tk.Button(button_frame,text="내보내기",command= lambda:conver_to_image_ppts(file_list.get(0,tk.END),organize_folder_ent.get()))
button_save.grid(row=4,column=0,sticky='we')

button_exit=tk.Button(button_frame,text="종료하기",command=root.destroy)
button_exit.grid(row=5,column=0,sticky='we')

file_frame_lbl['font']=button_font
button_load['font']=button_font
button_delete['font']=button_font
button_delete_all['font']=button_font
button_save['font']=button_font
button_exit['font']=button_font

#테스트용
if __name__=="__main__":
    organize_folder_ent['state'] = 'normal'
    organize_folder_ent.delete(0, tk.END)
    organize_folder_ent.insert(tk.END, 'C:/Users/HAN/Desktop/한국혁신연구원_자동화/PPTX')
root.mainloop()

# </editor-fold>

