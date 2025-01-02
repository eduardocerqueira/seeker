#date: 2025-01-02T16:59:29Z
#url: https://api.github.com/gists/2618cb81fe4ff45c77dec3c825ea4561
#owner: https://api.github.com/users/organisciak

#!/Users/peter.organisciak/anaconda3/bin python
# This script converts various filetypes to a markdown file, 
# saved in the same directory as the input file.

# Quick Action Instructions:
# It is useful with Mac OS Quick Actions, where you can select a file and 
# then select this script as the action.
# To set up a quick action, go to Automator, create a new quick action with the
# following settings:
# Workflow receives: files or folder in Finder
# Action is 'Run Shell Script'
# The shell script:
# for f in "$@"
# do
#     /path/to/your/python/python /path/to/your/script/convert_to_md.py "$f"
# done
import sys
from pathlib import Path
import re


def convert_pptx_to_md(input_path: Path, output_path: Path):
    from pptx import Presentation
    from markdownify import markdownify as md

    prs = Presentation(input_path)
    markdown_content = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                markdown_content += f"\n# {shape.text}\n"

    with open(output_path, "w") as md_file:
        md_file.write(md(markdown_content))


def convert_pdf_to_md(file_name: Path, output_file: Path):
    import PyPDF2

    markdown_content = ""
    with open(file_name, 'rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        for page_num, page in enumerate(pdf_reader.pages, 1):
            text = page.extract_text()
            markdown_content += f"# Page {page_num}\n\n{text}\n\n"
    
    with open(output_file, 'w', encoding='utf-8') as md_file:
        md_file.write(markdown_content)


def docx_markdown_formatting(paragraph):
    """
    This function receives a paragraph object and iterates over its runs,
    formatting bold and italic text with Markdown and returning the resulting string.
    """
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

    bold, bolditalic, italic, regular = 0, 0, 0, 0
    text = ""
    # count how much of the text is a certain style (for making heading inferences)
    for i, run in enumerate(paragraph.runs):
        text += run.text
        if run.text.strip() == '':
            continue
        if run.bold:
            bold += len(run.text)
        elif run.bold and run.italic:
            bolditalic += len(run.text)
        elif run.italic:
            italic += len(run.text)
        else:
            regular += len(run.text)

    total_len = len(text.strip())

    if paragraph.style.name == "Normal":
        # make guesses for a 'Normal formatted' paragraph
        if total_len <= bold:
            if paragraph.alignment == WD_PARAGRAPH_ALIGNMENT.CENTER:
                # Heading: Bold and Centered
                return f"# {text.strip()}"
            else:
                # Subheading: Bold
                return f"## {text.strip()}"
        elif total_len <= bolditalic:
            # Subsubheading: Bold and Italic
            return f"### {text.strip()}"
        elif total_len <= italic:
            # Subsubsubheading: Italic
            return f"#### {text.strip()}"
        else:
            # Apply formatting to individual runs
            formatted_text = ""
            for run in paragraph.runs:
                end_spaces = run.text[len(run.text.rstrip()):]
                start_spaces = run.text[:len(run.text) - len(run.text.lstrip())]
                if run.text.strip() == '':
                    formatted_text += run.text
                elif run.bold and run.italic:
                    # Both bold and italic
                    formatted_text += f"{start_spaces}***{run.text.strip()}***{end_spaces}"
                elif run.bold:
                    # Only bold
                    formatted_text += f"{start_spaces}**{run.text.strip()}**{end_spaces}"
                elif run.italic:
                    # Only italic
                    formatted_text += f"{start_spaces}_{run.text.strip()}_{end_spaces}"
                else:
                    # Regular text
                    formatted_text += run.text
            # fix any formatting within word
            formatted_text = formatted_text.replace('****', '').replace('’', "'").replace('“', '"').replace('”', '"')
            formatted_text = re.sub(r'\*\* +\*\*', r' ', formatted_text)
            # formatted_text = re.sub('(\W)([\*\_]{1,2})', r'\1 \2', formatted_text)
            # formatted_text = re.sub('(\w)[\*\_]{1,2}(\w)', r'\1\2', formatted_text)
            return formatted_text

    elif 'Heading' in paragraph.style.name:
        level = paragraph.style.name.split(' ')[1]  # Assume style is like 'Heading 1'
        return f"{'#' * int(level)} {text.strip()}"

    else:
        return text


def convert_docx_to_md(file_name: Path, output_file: Path):
    from docx import Document

    doc = Document(file_name)
    markdown_content = ""

    for para in doc.paragraphs:
        formatted_text = docx_markdown_formatting(para)
        markdown_content += formatted_text + "\n\n"

    with open(output_file, 'w', encoding='utf-8') as md_file:
        md_file.write(markdown_content)


def convert_by_file_name(file_name: Path):
    if not file_name.exists():
        print(f"File {file_name} does not exist")
        sys.exit(1)

    if file_name.is_dir():
        print(f"Input {file_name} is a directory")
        sys.exit(1)

    output_file = file_name.with_stem(file_name.stem + '_converted').with_suffix(".md")
    if file_name.suffix == ".pptx":
        convert_pptx_to_md(file_name, output_file)
    elif file_name.suffix == ".pdf":
        convert_pdf_to_md(file_name, output_file)
    elif file_name.suffix == ".docx":
        convert_docx_to_md(file_name, output_file)
    else:
        print(f"File {file_name} is not a pptx, pdf, or docx file")
        sys.exit(1)


if __name__ == "__main__":
    input_file = sys.argv[1]
    input_file = Path(input_file)
    convert_by_file_name(input_file)