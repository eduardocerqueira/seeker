#date: 2023-12-01T17:03:25Z
#url: https://api.github.com/gists/a07978387c02f313ee39be665b9d44eb
#owner: https://api.github.com/users/Mike3285

from pptx import Presentation
import os
import copy
def move_slide(copyFromPres: Presentation, slideIndex: int, pasteIntoPres: Presentation):
    """Takes two Presentation objs and an index, copies the slide with index slideIndex from copyFromPres to
    pasteIntoPres.
    returns the slide if everything went well, but the first presentation will contain all the other one's slides
    Thanks to https://stackoverflow.com/a/73954830/12380052 from which I copied the majority of this
    """
    # modeled on https://stackoverflow.com/a/56074651/20159015
    # and https://stackoverflow.com/a/62921848/20159015
    # take the slide at index slideIndex
    slide_to_copy = copyFromPres.slides[slideIndex]
    # Selecting the layout: it should and must be only one, so we take the 1st
    slide_layout = pasteIntoPres.slide_layouts[0]
    # names of other layouts can be found here under step 3:
    # https://www.geeksforgeeks.org/how-to-change-slide-layout-in-ms-powerpoint/
    # The layout we're using has an empty title with a placeholder like "Click to add title"

    # create now slide with that layout, to copy contents to
    new_slide = pasteIntoPres.slides.add_slide(slide_layout)
    # create dict for all the images it could find
    imgDict = {}  # entries will be generated if the pptx has images
    # variables for future use with hyperlinks
    haddress = None
    shape_w_hyperlink_id = None
    for shp in slide_to_copy.shapes:
        # Searching for images to not get a corrupt file in the end

        if 'Picture' in shp.name:
            # save image
            with open(shp.name + '.jpg', 'wb') as f:
                # Saving it temporarily
                f.write(shp.image.blob)
            # add image to dict
            imgDict[shp.name + '.jpg'] = [shp.left, shp.top, shp.width, shp.height]
        else:
            # create copy of elem
            el = shp.element
            text_frame = shp.text_frame
            p = text_frame.paragraphs[0]

            if p.runs:  # if this has ".runs" it's probably a hyperlink object
                # This
                r = p.runs[0]
                # let's save the ID of the shape with a hyperlink
                shape_w_hyperlink_id = shp.shape_id
                # ... and the hyperlink destination URI
                haddress = r.hyperlink.address

            newel = copy.deepcopy(el)
            # add elem to shape tree
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
            # If before we found a hyperlink...
        if haddress:
            for shape in new_slide.shapes:
                # let's check if we are in the shape with hyperlink
                if shape.shape_id == shape_w_hyperlink_id:
                    # we are on the shape with hyperlink, let's put it back where it belongs
                    # Not doing hyperlink.address=haddress because of the getter going on keyerror
                    shape.text_frame.paragraphs[0].runs[0].hyperlink._add_hlinkClick(haddress)

    # the following is from the guy on Stackoverflow:
    # things added first will be covered by things added last
    # => since I want pictures to be in foreground, I will add them after others elements
    # you can change this if you want to add pictures
    for k, v in imgDict.items():
        new_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])  # Adding the picture again
        os.remove(k)  # Removing the temp file it created
    try:
        new_slide.shapes.title.text = ' '  # todo it breaks if we delete this title box. We should find a way to delete it...
    except Exception as e:
        print("Warning, new_slide.shapes.title.text failed: %s" % e)
    return new_slide  # this returns the single slide so you can instantly work with it if you need to


def merge_presentations_list(paths: list, final_name: str = None):
    """mergio lista di presentazioni pptx in una unica"""
    outputPres = Presentation(paths[0])

    for path in paths[1:]:
        templatePres = Presentation(path)
        for i in range(len(templatePres.slides)):  # We iterate over all the slides
            move_slide(templatePres, i, outputPres)

    if final_name:
        outputPres.save(final_name)
    return outputPres


if __name__ == '__main__':
    presentation_list = [
        'pres1.pptx',
        'pres2.pptx',
        'with_link.pptx',
    ]
    merge_presentations_list(presentation_list, 'merged_presentations.pptx')
